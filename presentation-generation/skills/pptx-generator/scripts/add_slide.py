#!/usr/bin/env python3
"""Duplicate a slide from an existing PPTX or create from a layout.

Usage:
    python add_slide.py <input.pptx> --from-slide N --output <output.pptx>
    python add_slide.py <input.pptx> --from-layout N --output <output.pptx>

Options:
    --from-slide N    Duplicate slide N (1-based index)
    --from-layout N   Create new slide from layout N (1-based index)
    --output          Output file path
    --position N      Insert at position N (default: end)
"""

import sys
import os
import copy


def add_slide_from_existing(input_path, slide_idx, output_path, position=None):
    """Duplicate a slide in a PPTX file.

    Args:
        input_path: Path to input PPTX.
        slide_idx: 1-based index of slide to duplicate.
        output_path: Path for output PPTX.
        position: Insert position (1-based, default: end).

    Returns:
        tuple of (success: bool, error: str or None)
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError:
        return False, "python-pptx required. Install with: pip install python-pptx"

    if not os.path.isfile(input_path):
        return False, f"File not found: {input_path}"

    try:
        prs = Presentation(input_path)
        slides = prs.slides

        if slide_idx < 1 or slide_idx > len(slides):
            return False, f"Slide index {slide_idx} out of range (1-{len(slides)})"

        source_slide = slides[slide_idx - 1]

        # Use the same layout as the source slide
        layout = source_slide.slide_layout
        new_slide = prs.slides.add_slide(layout)

        # Copy shapes from source to new slide
        for shape in source_slide.shapes:
            el = copy.deepcopy(shape.element)
            new_slide.shapes._spTree.append(el)

        prs.save(output_path)
        return True, None
    except Exception as e:
        return False, str(e)


def add_slide_from_layout(input_path, layout_idx, output_path):
    """Add a new blank slide from a layout.

    Args:
        input_path: Path to input PPTX.
        layout_idx: 1-based index of slide layout.
        output_path: Path for output PPTX.

    Returns:
        tuple of (success: bool, error: str or None)
    """
    try:
        from pptx import Presentation
    except ImportError:
        return False, "python-pptx required. Install with: pip install python-pptx"

    try:
        prs = Presentation(input_path)
        layouts = prs.slide_layouts

        if layout_idx < 1 or layout_idx > len(layouts):
            return False, f"Layout index {layout_idx} out of range (1-{len(layouts)})"

        layout = layouts[layout_idx - 1]
        prs.slides.add_slide(layout)

        prs.save(output_path)
        return True, None
    except Exception as e:
        return False, str(e)


if __name__ == "__main__":
    if len(sys.argv) < 4:
        print(__doc__)
        sys.exit(1)

    input_path = sys.argv[1]
    output_path = None
    from_slide = None
    from_layout = None

    args = sys.argv[2:]
    i = 0
    while i < len(args):
        if args[i] == "--from-slide":
            from_slide = int(args[i + 1])
            i += 2
        elif args[i] == "--from-layout":
            from_layout = int(args[i + 1])
            i += 2
        elif args[i] == "--output":
            output_path = args[i + 1]
            i += 2
        else:
            i += 1

    if not output_path:
        print("Error: --output is required")
        sys.exit(1)

    if from_slide:
        success, error = add_slide_from_existing(input_path, from_slide, output_path)
    elif from_layout:
        success, error = add_slide_from_layout(input_path, from_layout, output_path)
    else:
        print("Error: specify --from-slide or --from-layout")
        sys.exit(1)

    if success:
        print(f"Slide added: {output_path}")
    else:
        print(f"Error: {error}")
        sys.exit(1)
